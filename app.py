import os
import base64
import subprocess
from flask import Flask, request, render_template, jsonify
from werkzeug.utils import secure_filename
import pdfplumber
from pptx import Presentation
import pandas as pd
import speech_recognition as sr
from sentence_transformers import SentenceTransformer
from groq import Groq
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import re
from datetime import datetime
from dotenv import load_dotenv  # Load from .env

# -------- Load environment variables securely --------
load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# -------- Groq client using API key from .env --------
groq_client = Groq(api_key=GROQ_API_KEY)
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')

# -------- Text Prettifier --------
def prettify_summary(text):
    text = text.strip()
    text = re.sub(r":([^\s])", r": \1", text)
    text = re.sub(r"(\n)?[-•*]\s", r"\n\n- ", text)
    text = re.sub(r"(?:^|\n)(#+\s?.+)", r"\n\n\1", text)
    text = re.sub(r"(?<!\n)\n(?!\n)", "\n\n", text)
    return text

# -------- File Extractor --------
def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.pdf':
        with pdfplumber.open(file_path) as pdf:
            return "\n".join([p.extract_text() or '' for p in pdf.pages])
    elif ext == '.pptx':
        prs = Presentation(file_path)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif ext == '.xlsx':
        df = pd.read_excel(file_path, sheet_name=None)
        return "\n".join([d.to_string() for d in df.values()])
    elif ext == '.csv':
        return pd.read_csv(file_path).to_string()
    elif ext == '.txt':
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    return ''

# -------- Summarization --------
def summarize_text(text):
    prompt = (
        "Summarize the following document. If there are tables, extract key insights. "
        "If there are images or charts, describe them if referenced. Present the summary clearly:\n\n"
        f"{text[:4000]}"
    )
    response = groq_client.chat.completions.create(
        model="llama3-8b-8192",
        messages=[{"role": "user", "content": prompt}]
    )
    return prettify_summary(response.choices[0].message.content.strip())

# -------- WAV Converter --------
def convert_to_wav(input_path, output_path):
    try:
        subprocess.run([
            'ffmpeg', '-y', '-i', input_path,
            '-ar', '16000', '-ac', '1', output_path
        ], check=True)
    except subprocess.CalledProcessError as e:
        print("FFmpeg error:", e)
        raise RuntimeError("Audio conversion failed")

# -------- Voice Transcription --------
def transcribe_audio(input_path):
    wav_path = input_path.replace('.webm', '.wav')
    convert_to_wav(input_path, wav_path)
    recognizer = sr.Recognizer()
    with sr.AudioFile(wav_path) as source:
        audio = recognizer.record(source)
    return recognizer.recognize_google(audio)

# -------- Live Audio Endpoint --------
@app.route('/voice-input', methods=['POST'])
def voice_input():
    audio_data = request.json.get("audio")
    if not audio_data:
        return jsonify({"error": "No audio received"}), 400

    audio_bytes = base64.b64decode(audio_data.split(",")[1])
    filename = f"voice_{datetime.now().strftime('%Y%m%d%H%M%S')}.webm"
    audio_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)

    with open(audio_path, 'wb') as f:
        f.write(audio_bytes)

    try:
        query = transcribe_audio(audio_path)
        return jsonify({"query": query})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# -------- Main Route --------
@app.route('/', methods=['GET', 'POST'])
def index():
    results = []
    query = ""

    if request.method == 'POST':
        uploaded_files = request.files.getlist("documents")
        query_text = request.form.get("query_text", "")

        file_texts, summaries, embeddings = {}, {}, {}

        for file in uploaded_files:
            filename = secure_filename(file.filename)
            path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(path)
            raw_text = extract_text(path)
            summary = summarize_text(raw_text)
            emb = embedding_model.encode(summary)
            file_texts[filename] = raw_text
            summaries[filename] = summary
            embeddings[filename] = emb

        query = query_text.strip()
        if query:
            query_emb = embedding_model.encode(query)

            dense_scores = {f: cosine_similarity([query_emb], [emb])[0][0] for f, emb in embeddings.items()}

            tfidf = TfidfVectorizer()
            corpus = list(summaries.values())
            X = tfidf.fit_transform(corpus)
            q_vec = tfidf.transform([query])
            sparse_scores = cosine_similarity(q_vec, X).flatten()

            final_scores = {}
            for i, fname in enumerate(summaries):
                final_scores[fname] = 0.5 * dense_scores[fname] + 0.5 * sparse_scores[i]

            results = sorted(
                [(fname, final_scores[fname], summaries[fname]) for fname in summaries],
                key=lambda x: x[1],
                reverse=True
            )

    return render_template('index.html', results=results, query=query)

if __name__ == '__main__':
    app.run(debug=True)
